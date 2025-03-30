import sys
import os
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
import pythoncom
import win32com.client

def convert_to_pptx(input_path, output_path):
    file_extension = os.path.splitext(input_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            # Convert PDF to PPTX
            # First convert PDF to images
            images = convert_from_path(input_path)
            
            # Create a new presentation
            prs = Presentation()
            
            # Add a slide for each page
            for image in images:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                
                # Save the image temporarily
                temp_image_path = f"temp_{len(prs.slides)}.png"
                image.save(temp_image_path)
                
                # Add the image to the slide
                left = 0
                top = 0
                pic = slide.shapes.add_picture(
                    temp_image_path, left, top, width=prs.slide_width, height=prs.slide_height
                )
                
                # Clean up temporary file
                os.remove(temp_image_path)
            
            # Save the presentation
            prs.save(output_path)
            
        elif file_extension in ['.jpg', '.jpeg', '.png']:
            # Convert image to PPTX
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            
            # Add the image to the slide
            left = 0
            top = 0
            pic = slide.shapes.add_picture(
                input_path, left, top, width=prs.slide_width, height=prs.slide_height
            )
            
            # Save the presentation
            prs.save(output_path)
            
        elif file_extension == '.docx':
            # Convert DOCX to PPTX using Word and PowerPoint
            pythoncom.CoInitialize()
            
            # Open Word document
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = True
            doc = word.Documents.Open(input_path)
            
            # Create PowerPoint presentation
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True
            presentation = powerpoint.Presentations.Add()
            
            # Process each paragraph as a new slide
            for paragraph in doc.Paragraphs:
                if paragraph.Range.Text.strip():
                    # Add new slide
                    slide = presentation.Slides.AddSlide(presentation.Slides.Count + 1, 1)
                    
                    # Add title
                    title_shape = slide.Shapes.Title
                    if title_shape:
                        title_shape.TextFrame.TextRange.Text = paragraph.Range.Text
            
            # Save and close
            presentation.SaveAs(output_path)
            presentation.Close()
            powerpoint.Quit()
            doc.Close()
            word.Quit()
            pythoncom.CoUninitialize()
        else:
            print(f"Unsupported file format: {file_extension}")
            sys.exit(1)
            
        print(f"Successfully converted {input_path} to PPTX")
        
    except Exception as e:
        print(f"Error converting file: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_to_pptx.py input_file output_file")
        sys.exit(1)
        
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    if not os.path.exists(input_path):
        print(f"Input file not found: {input_path}")
        sys.exit(1)
        
    convert_to_pptx(input_path, output_path) 